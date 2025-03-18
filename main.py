import streamlit as st
import pandas as pd
from pptx import Presentation
import os
import tempfile
from pathlib import Path
import subprocess

# Configuração da página
st.set_page_config(
    page_title="Gerador de estudos de viabilidade",
    page_icon="📄",
    layout="wide"
)

# Título da aplicação
st.title("Geração de estudo de viabilidade para implementação de Baterias")

# Função para carregar os dados do CSV
@st.cache_data
def load_data():
    try:
        df = pd.read_csv("relacao_empresas_cargas.csv", sep=",")
        # Renomeia as colunas para remover espaços
        df.columns = ["NOME_EMPRESARIAL", "NOME_CARGA","CNPJ_CARGA", "CODIGO_CARGA"]
        return df
    except Exception as e:
        st.error(f"Erro ao carregar o arquivo CSV: {str(e)}")
        return None

# Função para gerar o PDF
def generate_proposal(cliente, carga):
    try:
        # Verifica se o arquivo template existe
        template_path = "template_propostas.pptx"
        if not os.path.exists(template_path):
            st.error(f"Arquivo template não encontrado: {template_path}")
            return None
            
        # Carrega o template
        prs = Presentation(template_path)
        
        # Substitui os campos dinâmicos em todos os slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = shape.text.replace("{{NOME_CLIENTE}}", cliente)
                    shape.text = shape.text.replace("{{NOME_CARGA}}", carga)
        
        # Cria um diretório temporário para os arquivos
        with tempfile.TemporaryDirectory() as temp_dir:
            st.write(f"Diretório temporário criado: {temp_dir}")
            
            # Salva o PowerPoint modificado
            temp_pptx = os.path.join(temp_dir, "proposta_temp.pptx")
            prs.save(temp_pptx)
            st.write(f"PowerPoint temporário salvo em: {temp_pptx}")
            
            # Verifica se o LibreOffice está instalado
            if os.system("which soffice") != 0:
                st.error("LibreOffice não está instalado. Por favor, instale o LibreOffice para continuar.")
                return None
            
            # Converte para PDF usando LibreOffice (para Mac/Linux)
            temp_pdf = os.path.join(temp_dir, "proposta_temp.pdf")
            
            # Usa subprocess para capturar a saída do comando
            try:
                st.write("Iniciando conversão para PDF...")
                cmd = ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, temp_pptx]
                st.write(f"Comando a ser executado: {' '.join(cmd)}")
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True
                )
                
                st.write(f"Saída do comando: {result.stdout}")
                if result.stderr:
                    st.write(f"Erro do comando: {result.stderr}")
                
                if result.returncode != 0:
                    st.error(f"Erro ao converter PowerPoint para PDF. Código de retorno: {result.returncode}")
                    return None
                
                # Verifica se o PDF foi gerado
                if not os.path.exists(temp_pdf):
                    st.error(f"O arquivo PDF não foi gerado em: {temp_pdf}")
                    st.write("Conteúdo do diretório temporário:")
                    for file in os.listdir(temp_dir):
                        st.write(f"- {file}")
                    return None
                
                # Lê o PDF gerado
                with open(temp_pdf, "rb") as pdf_file:
                    pdf_bytes = pdf_file.read()
                
                return pdf_bytes
                
            except Exception as e:
                st.error(f"Erro durante a conversão: {str(e)}")
                return None
            
    except Exception as e:
        st.error(f"Erro ao gerar o estudo: {str(e)}")
        return None

# Carrega os dados
df = load_data()

if df is not None:
    # Cria duas colunas para os dropdowns
    col1, col2 = st.columns(2)
    
    with col1:
        # Dropdown para seleção do cliente
        clientes = sorted(df["NOME_EMPRESARIAL"].unique())
        cliente_selecionado = st.selectbox(
            "Selecione o Cliente",
            options=[""] + clientes,
            index=0
        )
    
    with col2:
        # Dropdown para seleção da carga
        if cliente_selecionado:
            cargas = sorted(df[df["NOME_EMPRESARIAL"] == cliente_selecionado]["NOME_CARGA"].unique())
            carga_selecionada = st.selectbox(
                "Selecione a Carga",
                options=[""] + cargas,
                index=0
            )
        else:
            carga_selecionada = st.selectbox(
                "Selecione a Carga",
                options=[""],
                index=0
            )
    
    # Botão para gerar a estudo
    if cliente_selecionado and carga_selecionada:
        if st.button("Gerar Estudo", type="primary"):
            with st.spinner("Gerando Estudo..."):
                pdf_bytes = generate_proposal(cliente_selecionado, carga_selecionada)
                
                if pdf_bytes:
                    # Cria o nome do arquivo
                    filename = f"estudo{cliente_selecionado.replace(' ', '_')}_{carga_selecionada.replace(' ', '_')}.pdf"
                    
                    # Botão para download
                    st.download_button(
                        label="Download do Estudo",
                        data=pdf_bytes,
                        file_name=filename,
                        mime="application/pdf"
                    )
                    st.success("Estudo gerado com sucesso!")
                else:
                    st.error("Erro ao gerar o estudo. Por favor, tente novamente.")
    else:
        st.info("Selecione um cliente e uma carga para gerar o estudo.")
else:
    st.error("Não foi possível carregar os dados. Verifique se o arquivo CSV está presente e no formato correto.") 